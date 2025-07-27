"""Document processing utilities for various file types."""

import os
from pathlib import Path
from typing import List, Dict, Any, Optional
from dataclasses import dataclass
import asyncio

import docx
import openpyxl
from pptx import Presentation
import PyPDF2
from PIL import Image
import pytesseract

from loguru import logger


@dataclass
class ProcessedDocument:
    """Represents a processed document with extracted content."""

    file_path: str
    file_type: str
    title: str
    content: str
    metadata: Dict[str, Any]
    chunks: List[str]


class DocumentProcessor:
    """Handles processing of various document types."""

    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    async def process_document(self, file_path: Path) -> Optional[ProcessedDocument]:
        """Process a single document and extract content."""
        try:
            file_extension = file_path.suffix.lower()

            if file_extension == ".txt":
                return await self._process_text_file(file_path)
            elif file_extension == ".docx":
                return await self._process_docx_file(file_path)
            elif file_extension in [".xlsx", ".xls"]:
                return await self._process_excel_file(file_path)
            elif file_extension in [".pptx", ".ppt"]:
                return await self._process_powerpoint_file(file_path)
            elif file_extension == ".pdf":
                return await self._process_pdf_file(file_path)
            elif file_extension in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
                return await self._process_image_file(file_path)
            else:
                logger.warning(f"Unsupported file type: {file_extension}")
                return None

        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            return None

    async def _process_text_file(self, file_path: Path) -> ProcessedDocument:
        """Process plain text files."""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()

        return ProcessedDocument(
            file_path=str(file_path),
            file_type="text",
            title=file_path.stem,
            content=content,
            metadata={"size": len(content)},
            chunks=self._chunk_text(content),
        )

    async def _process_docx_file(self, file_path: Path) -> ProcessedDocument:
        """Process Word documents."""
        doc = docx.Document(file_path)

        # Extract text from paragraphs
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        content = "\n".join(paragraphs)

        # Extract metadata
        metadata = {
            "author": doc.core_properties.author or "",
            "title": doc.core_properties.title or file_path.stem,
            "created": (
                str(doc.core_properties.created) if doc.core_properties.created else ""
            ),
            "modified": (
                str(doc.core_properties.modified)
                if doc.core_properties.modified
                else ""
            ),
            "paragraph_count": len(paragraphs),
        }

        return ProcessedDocument(
            file_path=str(file_path),
            file_type="docx",
            title=metadata["title"],
            content=content,
            metadata=metadata,
            chunks=self._chunk_text(content),
        )

    async def _process_excel_file(self, file_path: Path) -> ProcessedDocument:
        """Process Excel files."""
        workbook = openpyxl.load_workbook(file_path, data_only=True)

        content_parts = []
        sheet_info = {}

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_content = []

            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    row_text = "\t".join(
                        str(cell) if cell is not None else "" for cell in row
                    )
                    sheet_content.append(row_text)

            if sheet_content:
                content_parts.append(
                    f"Sheet: {sheet_name}\n" + "\n".join(sheet_content)
                )
                sheet_info[sheet_name] = len(sheet_content)

        content = "\n\n".join(content_parts)

        metadata = {
            "sheets": list(workbook.sheetnames),
            "sheet_info": sheet_info,
            "total_sheets": len(workbook.sheetnames),
        }

        return ProcessedDocument(
            file_path=str(file_path),
            file_type="excel",
            title=file_path.stem,
            content=content,
            metadata=metadata,
            chunks=self._chunk_text(content),
        )

    async def _process_powerpoint_file(self, file_path: Path) -> ProcessedDocument:
        """Process PowerPoint files."""
        prs = Presentation(file_path)

        content_parts = []
        slide_count = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())

            if slide_content:
                content_parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_content))
                slide_count += 1

        content = "\n\n".join(content_parts)

        metadata = {"slide_count": slide_count, "total_slides": len(prs.slides)}

        return ProcessedDocument(
            file_path=str(file_path),
            file_type="powerpoint",
            title=file_path.stem,
            content=content,
            metadata=metadata,
            chunks=self._chunk_text(content),
        )

    async def _process_pdf_file(self, file_path: Path) -> ProcessedDocument:
        """Process PDF files."""
        content_parts = []

        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)

            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    text = page.extract_text()
                    if text.strip():
                        content_parts.append(f"Page {page_num}:\n{text}")
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num}: {e}")

        content = "\n\n".join(content_parts)

        metadata = {"page_count": len(pdf_reader.pages)}

        return ProcessedDocument(
            file_path=str(file_path),
            file_type="pdf",
            title=file_path.stem,
            content=content,
            metadata=metadata,
            chunks=self._chunk_text(content),
        )

    async def _process_image_file(self, file_path: Path) -> ProcessedDocument:
        """Process image files using OCR."""
        try:
            image = Image.open(file_path)
            content = pytesseract.image_to_string(image)

            metadata = {"size": image.size, "mode": image.mode, "format": image.format}

            return ProcessedDocument(
                file_path=str(file_path),
                file_type="image",
                title=file_path.stem,
                content=content,
                metadata=metadata,
                chunks=self._chunk_text(content) if content.strip() else [],
            )
        except Exception as e:
            logger.error(f"Error processing image {file_path}: {e}")
            return None

    def _chunk_text(self, text: str) -> List[str]:
        """Split text into overlapping chunks."""
        if not text.strip():
            return []

        chunks = []
        start = 0

        while start < len(text):
            end = start + self.chunk_size

            # Try to break at sentence or paragraph boundaries
            if end < len(text):
                # Look for sentence endings
                for i in range(end, max(start + self.chunk_size // 2, end - 100), -1):
                    if text[i] in ".!?\n":
                        end = i + 1
                        break

            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)

            start = end - self.chunk_overlap
            if start >= len(text):
                break

        return chunks
